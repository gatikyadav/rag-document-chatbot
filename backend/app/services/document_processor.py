import os
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup
import markdown
from datetime import datetime
import hashlib

from ..config import settings
from ..models.schemas import DocumentChunk, DocumentMetadata

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles processing of various document types into text chunks."""
    
    def __init__(self, chunk_size: int = None, chunk_overlap: int = None):
        self.chunk_size = chunk_size or settings.chunk_size
        self.chunk_overlap = chunk_overlap or settings.chunk_overlap
        
        # Supported file extensions
        self.supported_extensions = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.pptx': self._extract_pptx,
            '.xlsx': self._extract_excel,
            '.xls': self._extract_excel,
            '.txt': self._extract_txt,
            '.md': self._extract_markdown,
            '.html': self._extract_html
        }
        
        logger.info(f"DocumentProcessor initialized with chunk_size={self.chunk_size}, chunk_overlap={self.chunk_overlap}")
    
    def process_file(self, file_path: str) -> List[DocumentChunk]:
        """Process a file and return document chunks with metadata."""
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if file_path.suffix.lower() not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_path.suffix}")
        
        logger.info(f"Processing file: {file_path.name}")
        
        try:
            # Extract text based on file type
            text_content = self._extract_text(file_path)
            
            if not text_content.strip():
                logger.warning(f"No text extracted from {file_path}")
                return []
            
            # Create chunks
            text_chunks = self._create_chunks(text_content)
            
            # Convert to DocumentChunk objects with metadata
            document_chunks = []
            for i, chunk_text in enumerate(text_chunks):
                metadata = self._create_metadata(file_path, i, len(text_chunks))
                
                chunk = DocumentChunk(
                    text=chunk_text,
                    metadata=metadata
                )
                document_chunks.append(chunk)
            
            logger.info(f"Successfully processed {file_path.name}: {len(document_chunks)} chunks created")
            return document_chunks
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")
            raise
    
    def process_directory(self, directory_path: str) -> List[DocumentChunk]:
        """Process all supported files in a directory."""
        directory = Path(directory_path)
        
        if not directory.exists():
            raise FileNotFoundError(f"Directory not found: {directory_path}")
        
        all_chunks = []
        processed_files = 0
        failed_files = []
        
        # Find all supported files
        for file_path in directory.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_extensions:
                try:
                    chunks = self.process_file(str(file_path))
                    all_chunks.extend(chunks)
                    processed_files += 1
                except Exception as e:
                    logger.error(f"Failed to process {file_path}: {str(e)}")
                    failed_files.append(str(file_path))
        
        logger.info(f"Directory processing complete: {processed_files} files processed, {len(failed_files)} failed, {len(all_chunks)} total chunks")
        
        if failed_files:
            logger.warning(f"Failed files: {failed_files}")
        
        return all_chunks
    
    def _extract_text(self, file_path: Path) -> str:
        """Extract text from different file types."""
        suffix = file_path.suffix.lower()
        
        if suffix in self.supported_extensions:
            return self.supported_extensions[suffix](file_path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")
    
    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {str(e)}")
            raise
        return text
    
    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX files."""
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text += row_text + "\n"
                        
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {str(e)}")
            raise
        return text
    
    def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX files."""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {str(e)}")
            raise
        return text
    
    def _extract_excel(self, file_path: Path) -> str:
        """Extract text from Excel files."""
        text = ""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip() and row_text != " | " * (len(row) - 1):
                        text += row_text + "\n"
        except Exception as e:
            logger.error(f"Error extracting Excel {file_path}: {str(e)}")
            raise
        return text
    
    def _extract_txt(self, file_path: Path) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error extracting TXT {file_path}: {str(e)}")
            raise
    
    def _extract_markdown(self, file_path: Path) -> str:
        """Extract text from Markdown files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                md_content = file.read()
                # Convert markdown to HTML then extract text
                html = markdown.markdown(md_content)
                soup = BeautifulSoup(html, 'html.parser')
                return soup.get_text()
        except Exception as e:
            logger.error(f"Error extracting Markdown {file_path}: {str(e)}")
            raise
    
    def _extract_html(self, file_path: Path) -> str:
        """Extract text from HTML files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                return soup.get_text()
        except Exception as e:
            logger.error(f"Error extracting HTML {file_path}: {str(e)}")
            raise
    
    def _create_chunks(self, text: str) -> List[str]:
        """Split text into overlapping chunks."""
        if not text.strip():
            return []
        
        # Split by words for better chunking
        words = text.split()
        
        if len(words) <= self.chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(words):
            end = start + self.chunk_size
            chunk_words = words[start:end]
            chunk_text = " ".join(chunk_words)
            chunks.append(chunk_text)
            
            # Move start position considering overlap
            start += self.chunk_size - self.chunk_overlap
            
            # Ensure we don't go beyond the text
            if start >= len(words):
                break
        
        return chunks
    
    def _create_metadata(self, file_path: Path, chunk_id: int, total_chunks: int) -> DocumentMetadata:
        """Create metadata for a document chunk."""
        # Convert to absolute path to handle relative paths properly
        abs_file_path = file_path.resolve()
    
        # Generate public URL for the document
        public_url = f"/static/documents/{file_path.name}"
    
        # Get file stats
        file_stats = abs_file_path.stat()
        file_size = self._format_file_size(file_stats.st_size)
        created_date = datetime.fromtimestamp(file_stats.st_mtime)
    
        return DocumentMetadata(
            source_file=str(abs_file_path),
            filename=file_path.name,
            file_type=file_path.suffix.lower().replace('.', ''),
            chunk_id=chunk_id,
            public_url=public_url,
            file_size=file_size,
            created_date=created_date
        )
    
    def _format_file_size(self, size_bytes: int) -> str:
        """Format file size in human readable format."""
        if size_bytes == 0:
            return "0 B"
        
        size_names = ["B", "KB", "MB", "GB"]
        i = 0
        while size_bytes >= 1024 and i < len(size_names) - 1:
            size_bytes /= 1024.0
            i += 1
        
        return f"{size_bytes:.1f} {size_names[i]}"
    
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return list(self.supported_extensions.keys())
    
    def is_supported_file(self, file_path: str) -> bool:
        """Check if a file type is supported."""
        return Path(file_path).suffix.lower() in self.supported_extensions